# src/indexer.py

import os
import time
import uuid
import sqlite3
from pathlib import Path
from datetime import datetime

from tqdm import tqdm
from qdrant_client import QdrantClient
from qdrant_client.http.models import Distance, VectorParams, PointStruct
from sentence_transformers import SentenceTransformer

import fitz  # PyMuPDF
from PIL import Image
import pytesseract

import docx
import pptx
import logging

# -----------------------------
# Paths & Config
# -----------------------------
BASE_DIR        = Path(__file__).parent.parent.resolve()
#DATA_DIR        = BASE_DIR / "data" #commenting this out for a bit
DATA_DIR = BASE_DIR / "data" / "test"  # For testing
LOGS_DIR        = BASE_DIR / "logs"
SQLITE_DB       = BASE_DIR / "data" / "index_metadata.db"

COLLECTION_NAME = "optimAIze-index"
QDRANT_URL      = "http://localhost:6333"

# Ensure directories exist
DATA_DIR.mkdir(exist_ok=True)
LOGS_DIR.mkdir(exist_ok=True)

# -----------------------------
# Logging Setup
# -----------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.FileHandler(LOGS_DIR / "indexer.log"),
        logging.StreamHandler()  # Also log to console
    ]
)

logger = logging.getLogger(__name__)

# -----------------------------
# Initialize SQLite Metadata DB
# -----------------------------
def init_metadata_db():
    """Initialize SQLite database with proper error handling"""
    try:
        conn = sqlite3.connect(SQLITE_DB)
        c = conn.cursor()
        
        # Table for document‐level metadata
        c.execute("""
          CREATE TABLE IF NOT EXISTS documents (
            document_id   TEXT PRIMARY KEY,
            file_path     TEXT UNIQUE NOT NULL,
            file_name     TEXT NOT NULL,
            title         TEXT,
            author        TEXT,
            last_modified REAL,
            category      TEXT
          )
        """)
        
        # Table for chunk‐level text
        c.execute("""
          CREATE TABLE IF NOT EXISTS chunks (
            chunk_id     TEXT PRIMARY KEY,
            document_id  TEXT NOT NULL,
            chunk_text   TEXT NOT NULL,
            FOREIGN KEY(document_id) REFERENCES documents(document_id)
          )
        """)
        
        # FTS5 index on chunk_text (for keyword search)
        c.execute("""
          CREATE VIRTUAL TABLE IF NOT EXISTS chunks_fts
          USING fts5(
            chunk_text,
            content='chunks',
            content_rowid='rowid'
          )
        """)
        
        conn.commit()
        logger.info("SQLite database initialized successfully")
        return conn
        
    except Exception as e:
        logger.error(f"Failed to initialize SQLite database: {e}")
        raise

# -----------------------------
# Load Embedding Model
# -----------------------------
def load_embedding_model():
    """Load embedding model with proper error handling"""
    try:
        logger.info("⏳ Loading local embedding model (nomic-embed-text-v1)...")
        model = SentenceTransformer(
            "nomic-ai/nomic-embed-text-v1",
            trust_remote_code=True
        )
        embed_dim = model.get_sentence_embedding_dimension()
        logger.info(f"✅ Model loaded successfully (dim = {embed_dim})")
        return model, embed_dim
    except Exception as e:
        logger.error(f"Failed to load embedding model: {e}")
        raise

# -----------------------------
# Qdrant Setup
# -----------------------------
def setup_qdrant_client(embed_dim):
    """Setup Qdrant client with proper error handling"""
    try:
        client = QdrantClient(url=QDRANT_URL)
        
        # Test connection
        collections = client.get_collections()
        logger.info("✅ Connected to Qdrant successfully")
        
        # Check if collection exists, recreate if needed
        try:
            collection_info = client.get_collection(COLLECTION_NAME)
            logger.info(f"Collection '{COLLECTION_NAME}' already exists")
        except Exception:
            logger.info(f"Creating new collection '{COLLECTION_NAME}'")
            client.recreate_collection(
                collection_name=COLLECTION_NAME,
                vectors_config=VectorParams(size=embed_dim, distance=Distance.COSINE),
            )
        
        return client
        
    except Exception as e:
        logger.error(f"Failed to setup Qdrant client: {e}")
        logger.error("Make sure Qdrant is running: docker run -p 6333:6333 qdrant/qdrant")
        raise

# -----------------------------
# File Parsers
# -----------------------------
def parse_txt(file_path: Path) -> str:
    """Parse text file with better error handling"""
    try:
        text = file_path.read_text(encoding="utf-8", errors="ignore")
        return text.strip()
    except Exception as e:
        logger.warning(f"Failed to parse TXT {file_path}: {e}")
        return ""

def parse_pdf(file_path: Path) -> str:
    """Parse PDF with OCR fallback and better error handling"""
    try:
        doc = fitz.open(str(file_path))
        full_text = []
        
        for page_num, page in enumerate(doc):
            try:
                text = page.get_text().strip()
                if text:
                    full_text.append(text)
                else:
                    # No embedded text → render page as image and OCR
                    try:
                        pix = page.get_pixmap(dpi=200)
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        ocr_text = pytesseract.image_to_string(img).strip()
                        if ocr_text:
                            full_text.append(ocr_text)
                    except Exception as ocr_e:
                        logger.warning(f"OCR failed for page {page_num} of {file_path}: {ocr_e}")
                        
            except Exception as page_e:
                logger.warning(f"Failed to process page {page_num} of {file_path}: {page_e}")
                
        doc.close()
        return "\n".join(full_text)
        
    except Exception as e:
        logger.warning(f"Failed to parse PDF {file_path}: {e}")
        return ""

def parse_docx(file_path: Path) -> str:
    """Parse DOCX with better error handling"""
    try:
        doc = docx.Document(str(file_path))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except Exception as e:
        logger.warning(f"Failed to parse DOCX {file_path}: {e}")
        return ""

def parse_pptx(file_path: Path) -> str:
    """Parse PPTX with better error handling"""
    try:
        prs = pptx.Presentation(str(file_path))
        text = []
        for slide_num, slide in enumerate(prs.slides):
            try:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text.strip())
            except Exception as slide_e:
                logger.warning(f"Failed to process slide {slide_num} of {file_path}: {slide_e}")
        return "\n".join(text)
    except Exception as e:
        logger.warning(f"Failed to parse PPTX {file_path}: {e}")
        return ""

# Map extensions → parser functions
SUPPORTED_EXTS = {
    ".txt": parse_txt,
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
}

# -----------------------------
# Chunking Function
# -----------------------------
def chunk_text(text: str, max_len: int = 500, overlap: int = 50) -> list[str]:
    """
    Improved chunking with overlap and sentence boundary awareness
    """

    print("DEBUG: chunk_text function called with text length:", len(text) if text else 0)  # ADD THIS LINE

    if not text or not text.strip():
        return []
    
    # Clean up the text
    text = " ".join(text.split())  # Normalize whitespace
    
    if len(text) <= max_len:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + max_len
        
        if end >= len(text):
            # Last chunk
            chunk = text[start:].strip()
            if chunk:
                chunks.append(chunk)
            break
        
        # Try to find a good breaking point (sentence end, then word boundary)
        chunk_content = text[start:end]  # RENAMED from chunk_text to chunk_content
        
        # Look for sentence endings
        last_sentence = max(
            chunk_content.rfind('. '),
            chunk_content.rfind('! '),
            chunk_content.rfind('? ')
        )
        
        if last_sentence > len(chunk_content) * 0.5:  # Good sentence break found
            end = start + last_sentence + 1
        else:
            # Look for word boundary
            last_space = chunk_content.rfind(' ')
            if last_space > len(chunk_content) * 0.7:  # Good word break found
                end = start + last_space
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        # Move start position with overlap
        start = end - overlap
        if start <= 0:
            start = end
    
    return [chunk for chunk in chunks if chunk.strip()]

# -----------------------------
# Insert or Retrieve Document Metadata
# -----------------------------
def upsert_document_metadata(conn: sqlite3.Connection, filepath: Path) -> str:
    """
    If the file_path is already in the documents table, return its document_id.
    Otherwise, insert a new row and return the newly generated document_id.
    """
    try:
        c = conn.cursor()
        file_path_str = str(filepath.resolve())
        
        # Check if already exists
        c.execute("SELECT document_id FROM documents WHERE file_path = ?", (file_path_str,))
        row = c.fetchone()
        if row:
            return row[0]

        # Not present → insert
        document_id = str(uuid.uuid4())
        file_name     = filepath.name
        title         = filepath.stem
        author        = None
        last_modified = filepath.stat().st_mtime
        category      = filepath.parent.name

        c.execute("""
          INSERT INTO documents (document_id, file_path, file_name, title, author, last_modified, category)
          VALUES (?, ?, ?, ?, ?, ?, ?)
        """, (
            document_id,
            file_path_str,
            file_name,
            title,
            author,
            last_modified,
            category
        ))
        conn.commit()
        return document_id
        
    except Exception as e:
        logger.error(f"Failed to upsert document metadata for {filepath}: {e}")
        raise

# -----------------------------
# Insert One Chunk Into SQLite + FTS
# -----------------------------
def insert_chunk_metadata(
    conn: sqlite3.Connection,
    document_id: str,
    chunk_content: str  # CHANGED from chunk_text
) -> str:
    """
    Inserts a chunk into the `chunks` table and its text into `chunks_fts`.
    Returns the generated chunk_id.
    """
    try:
        c = conn.cursor()
        chunk_id = str(uuid.uuid4())
        
        # Insert into chunks table
        c.execute("""
          INSERT INTO chunks (chunk_id, document_id, chunk_text)
          VALUES (?, ?, ?)
        """, (chunk_id, document_id, chunk_content))  # CHANGED
        
        rowid = c.lastrowid

        # Insert into FTS5 index
        c.execute("INSERT INTO chunks_fts(rowid, chunk_text) VALUES (?, ?)", (rowid, chunk_content))  # CHANGED
        conn.commit()
        return chunk_id
        
    except Exception as e:
        logger.error(f"Failed to insert chunk metadata: {e}")
        raise

# -----------------------------
# Process & Index One File
# -----------------------------
def process_and_index_file(
    conn: sqlite3.Connection,
    filepath: Path,
    client: QdrantClient,
    model: SentenceTransformer
):
    """
    Process and index a single file with comprehensive error handling
    """
    try:
        logger.info(f"Processing: {filepath}")

        # DEBUG: Check file size
        file_size = filepath.stat().st_size
        logger.info(f"File size: {file_size} bytes")

        # 1) Parse file → text
        logger.info("Step 1: Parsing file...")
        parser = SUPPORTED_EXTS[filepath.suffix.lower()]
        text = parser(filepath)

        logger.info(f"Step 1 complete: extracted {len(text) if text else 0} characters")

        if not text or not text.strip():
            logger.warning(f"Skipped unreadable/empty: {filepath}")
            return

        # 2) Chunk text
        logger.info("Step 2: Chunking text...")
        chunks = chunk_text(text, max_len=500)
        if not chunks:
            logger.warning(f"No chunks returned for: {filepath}")
            return

        logger.info(f"Step 2 complete: generated {len(chunks)} chunks")

        # 3) Upsert document metadata
        document_id = upsert_document_metadata(conn, filepath)

        # 4) Insert each chunk into SQLite
        chunk_ids = []
        for i, chunk in enumerate(chunks):
            try:
                chunk_id = insert_chunk_metadata(conn, document_id, chunk)
                chunk_ids.append(chunk_id)
            except Exception as e:
                logger.error(f"Failed to insert chunk {i} for {filepath}: {e}")
                continue

        if not chunk_ids:
            logger.warning(f"No chunks successfully inserted for {filepath}")
            return

        # 5) Generate embeddings
        logger.info(f"Generating embeddings for {len(chunks)} chunks...")
        try:
            embeddings = model.encode(chunks, show_progress_bar=False)
            embeddings = embeddings.tolist()
        except Exception as e:
            logger.error(f"Failed to generate embeddings for {filepath}: {e}")
            return

        # 6) Prepare Qdrant points
        points = []
        timestamp = time.time()
        
        for i, (chunk_id, chunk_content, emb) in enumerate(zip(chunk_ids, chunks, embeddings)):
            try:
                payload = {
                    "file_path":     str(filepath.resolve()),
                    "chunk":         chunk_content[:500],  # Truncate for payload size
                    "timestamp":     timestamp,
                    "file_type":     filepath.suffix.lower(),
                    "document_id":   document_id,
                }
                points.append(
                    PointStruct(
                        id=chunk_id, 
                        vector=emb, 
                        payload=payload
                    )
                )
            except Exception as e:
                logger.error(f"Failed to create point {i} for {filepath}: {e}")
                continue

        # 7) Upsert to Qdrant
        if points:
            try:
                client.upsert(collection_name=COLLECTION_NAME, points=points)
                logger.info(f"Successfully indexed: {filepath} ({len(points)} chunks)")
            except Exception as e:
                logger.error(f"Failed to upsert to Qdrant for {filepath}: {e}")
                return
        else:
            logger.warning(f"No valid points created for {filepath}")
            
    except Exception as e:
        logger.error(f"Failed to process {filepath}: {e}")

# -----------------------------
# Walk Directory & Index All Files
# -----------------------------
def index_directory(
    conn: sqlite3.Connection,
    data_path: Path,
    client: QdrantClient,
    model: SentenceTransformer
):
    """
    Recursively walk `data_path` and index all supported files
    """
    # Find all supported files
    supported_files = []
    for filepath in data_path.rglob("*"):
        if not filepath.is_file(): 
            continue
        if filepath.suffix.lower() not in SUPPORTED_EXTS:
            continue
        # Skip database file
        if filepath.name == "index_metadata.db":
            continue
        supported_files.append(filepath)
    
    if not supported_files:
        logger.warning(f"No supported files found in {data_path}")
        logger.info(f"Supported extensions: {list(SUPPORTED_EXTS.keys())}")
        return
    
    logger.info(f"Found {len(supported_files)} supported files to index")
    
    # Process each file
    for filepath in tqdm(supported_files, desc="Indexing files"):
        try:
            process_and_index_file(conn, filepath, client, model)
        except Exception as e:
            logger.error(f"Critical error processing {filepath}: {e}")
            continue

# -----------------------------
# Main
# -----------------------------
def main():
    """Main function with comprehensive error handling"""
    try:
        logger.info("🔍 Starting OptimAIze indexer...")
        logger.info(f"Data directory: {DATA_DIR}")
        logger.info(f"SQLite database: {SQLITE_DB}")
        logger.info(f"Qdrant URL: {QDRANT_URL}")
        
        # Check if data directory exists and has files
        if not DATA_DIR.exists():
            logger.error(f"Data directory does not exist: {DATA_DIR}")
            return
        
        # Initialize components
        logger.info("Initializing SQLite database...")
        conn = init_metadata_db()
        
        logger.info("Loading embedding model...")
        model, embed_dim = load_embedding_model()
        
        logger.info("Setting up Qdrant client...")
        client = setup_qdrant_client(embed_dim)
        
        # Index files
        logger.info("Starting file indexing...")
        index_directory(conn, DATA_DIR, client, model)
        
        # Cleanup
        conn.close()
        logger.info("✅ Indexing completed successfully")
        
    except KeyboardInterrupt:
        logger.info("Indexing interrupted by user")
    except Exception as e:
        logger.error(f"Critical error in main: {e}")
        raise

if __name__ == "__main__":
    main()