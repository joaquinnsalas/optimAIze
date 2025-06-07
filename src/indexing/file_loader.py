"""File loading utilities for OptimAIze."""

import io
from pathlib import Path
from typing import List, Dict, Any, Optional
import PyPDF2
import pytesseract
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
import openpyxl
from PIL import Image
from src.config.settings import config
from src.utils.logger import logger

class FileLoader:
    """File loader for various document types."""
    
    def __init__(self):
        self.ocr_dpi = config.indexing.get("ocr_dpi", 300)
    
    def load_file(self, file_path: Path) -> Optional[Dict[str, Any]]:
        """Load content from a file based on its extension."""
        try:
            extension = file_path.suffix.lower()
            
            if extension == '.pdf':
                return self._load_pdf(file_path)
            elif extension == '.docx':
                return self._load_docx(file_path)
            elif extension == '.pptx':
                return self._load_pptx(file_path)
            elif extension == '.txt':
                return self._load_txt(file_path)
            elif extension == '.md':
                return self._load_md(file_path)
            elif extension == '.xlsx':
                return self._load_xlsx(file_path)
            else:
                logger.warning(f"Unsupported file extension: {extension}")
                return None
        
        except Exception as e:
            logger.error(f"Error loading PPTX {file_path}: {e}")
            return {"content": "", "metadata": {"source": str(file_path), "type": "pptx"}}
    
    def _load_txt(self, file_path: Path) -> Dict[str, Any]:
        """Load TXT content."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            metadata = {
                "source": str(file_path),
                "type": "txt",
                "encoding": "utf-8"
            }
            
            return {"content": content, "metadata": metadata}
        
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    content = file.read()
                metadata = {
                    "source": str(file_path),
                    "type": "txt",
                    "encoding": "latin-1"
                }
                return {"content": content, "metadata": metadata}
            except Exception as e:
                logger.error(f"Error loading TXT {file_path}: {e}")
                return {"content": "", "metadata": {"source": str(file_path), "type": "txt"}}
        
        except Exception as e:
            logger.error(f"Error loading TXT {file_path}: {e}")
            return {"content": "", "metadata": {"source": str(file_path), "type": "txt"}}
    
    def _load_md(self, file_path: Path) -> Dict[str, Any]:
        """Load Markdown content."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            metadata = {
                "source": str(file_path),
                "type": "markdown",
                "encoding": "utf-8"
            }
            
            return {"content": content, "metadata": metadata}
        
        except Exception as e:
            logger.error(f"Error loading Markdown {file_path}: {e}")
            return {"content": "", "metadata": {"source": str(file_path), "type": "markdown"}}
    
    def _load_xlsx(self, file_path: Path) -> Dict[str, Any]:
        """Load XLSX content."""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            content = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content += f"\n--- Sheet: {sheet_name} ---\n"
                
                # Get all rows with data
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        content += row_text + "\n"
            
            metadata = {
                "source": str(file_path),
                "type": "xlsx",
                "sheets": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames
            }
            
            return {"content": content, "metadata": metadata}
        
        except Exception as e:
            logger.error(f"Error loading XLSX {file_path}: {e}")
            return {"content": "", "metadata": {"source": str(file_path), "type": "xlsx"}}

    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions."""
        return ['.pdf', '.docx', '.pptx', '.txt', '.md', '.xlsx']
    
    def validate_file(self, file_path: Path) -> bool:
        """Validate if file can be processed."""
        try:
            if not file_path.exists():
                logger.warning(f"File does not exist: {file_path}")
                return False
            
            if not file_path.is_file():
                logger.warning(f"Path is not a file: {file_path}")
                return False
            
            extension = file_path.suffix.lower()
            if extension not in self.get_supported_extensions():
                logger.warning(f"Unsupported file extension: {extension}")
                return False
            
            # Check file size (warn if very large)
            file_size = file_path.stat().st_size
            if file_size > 100 * 1024 * 1024:  # 100MB
                logger.warning(f"Large file detected: {file_path} ({file_size / 1024 / 1024:.1f}MB)")
            
            return True
        
        except Exception as e:
            logger.error(f"Error validating file {file_path}: {e}")
            return False
    
    def get_file_info(self, file_path: Path) -> Dict[str, Any]:
        """Get basic information about a file without loading content."""
        try:
            if not self.validate_file(file_path):
                return {}
            
            stat = file_path.stat()
            extension = file_path.suffix.lower()
            
            info = {
                "path": str(file_path),
                "name": file_path.name,
                "extension": extension,
                "size": stat.st_size,
                "size_human": self._format_file_size(stat.st_size),
                "modified_time": stat.st_mtime,
                "can_process": True
            }
            
            # Add format-specific info
            if extension == '.pdf':
                info["supports_ocr"] = True
            elif extension in ['.docx', '.pptx']:
                info["office_format"] = True
            elif extension == '.xlsx':
                info["spreadsheet"] = True
            
            return info
        
        except Exception as e:
            logger.error(f"Error getting file info for {file_path}: {e}")
            return {"path": str(file_path), "can_process": False, "error": str(e)}
    
    def _format_file_size(self, size_bytes: int) -> str:
        """Format file size in human readable format."""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024.0:
                return f"{size_bytes:.1f} {unit}"
            size_bytes /= 1024.0
        return f"{size_bytes:.1f} TB"
    
    def batch_load_files(self, file_paths: List[Path]) -> List[Dict[str, Any]]:
        """Load multiple files with progress tracking."""
        documents = []
        
        logger.info(f"Starting batch load of {len(file_paths)} files")
        
        for i, file_path in enumerate(file_paths):
            try:
                logger.debug(f"Loading file {i+1}/{len(file_paths)}: {file_path}")
                
                document = self.load_file(file_path)
                if document:
                    documents.append(document)
                else:
                    logger.warning(f"Failed to load content from {file_path}")
            
            except Exception as e:
                logger.error(f"Error in batch loading file {file_path}: {e}")
                continue
        
        logger.info(f"Successfully loaded {len(documents)} out of {len(file_paths)} files")
        return documents
    
    def _load_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Load PDF content with OCR fallback for image-based PDFs."""
        content = ""
        metadata = {"source": str(file_path), "type": "pdf", "pages": 0}
        
        try:
            # First try to extract text directly
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                metadata["pages"] = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        content += f"\n--- Page {page_num + 1} ---\n{text}"
            
            # If no text extracted or very little text, use OCR
            if len(content.strip()) < 100:
                logger.info(f"PDF {file_path} appears to be image-based, using OCR")
                content = self._ocr_pdf(file_path)
                metadata["ocr_used"] = True
            else:
                metadata["ocr_used"] = False
            
            return {"content": content, "metadata": metadata}
        
        except Exception as e:
            logger.error(f"Error loading PDF {file_path}: {e}")
            return {"content": "", "metadata": metadata}
    
    def _ocr_pdf(self, file_path: Path) -> str:
        """Extract text from PDF using OCR."""
        try:
            # Convert PDF to images
            images = convert_from_path(file_path, dpi=self.ocr_dpi)
            content = ""
            
            for i, image in enumerate(images):
                # Extract text using OCR
                text = pytesseract.image_to_string(image, lang='eng')
                if text.strip():
                    content += f"\n--- Page {i + 1} (OCR) ---\n{text}"
            
            return content
        
        except Exception as e:
            logger.error(f"Error performing OCR on PDF {file_path}: {e}")
            return ""
    
    def _load_docx(self, file_path: Path) -> Dict[str, Any]:
        """Load DOCX content."""
        try:
            doc = Document(file_path)
            content = ""
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content += paragraph.text + "\n"
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        content += row_text + "\n"
            
            metadata = {
                "source": str(file_path),
                "type": "docx",
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables)
            }
            
            return {"content": content, "metadata": metadata}
        
        except Exception as e:
            logger.error(f"Error loading DOCX {file_path}: {e}")
            return {"content": "", "metadata": {"source": str(file_path), "type": "docx"}}
    
    def _load_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Load PPTX content."""
        try:
            prs = Presentation(file_path)
            content = ""
            
            for slide_num, slide in enumerate(prs.slides):
                slide_content = f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content += shape.text + "\n"
                
                if slide_content.strip() != f"--- Slide {slide_num + 1} ---":
                    content += slide_content
            
            metadata = {
                "source": str(file_path),
                "type": "pptx",
                "slides": len(prs.slides)
            }
            
            return {"content": content, "metadata": metadata}
        
        except Exception as e:
            logger.error(f"Error loading PPTX {file_path}: {e}")
            return {"content": "", "metadata": {"source": str(file_path), "type": "pptx"}}