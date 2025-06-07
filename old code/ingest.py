# src/ingest.py - Fixed version with schema migration and better error handling
import logging
import sqlite3
from pathlib import Path
from typing import List, Dict, Optional, Any, Tuple
import uuid
from datetime import datetime
import hashlib
import os
import sys
import time
import json

# Document processing imports
import fitz  # PyMuPDF
from PIL import Image
import pytesseract
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
import csv

# Vector store and embeddings
from qdrant_client import QdrantClient
from qdrant_client.http.models import Distance, VectorParams, PointStruct
from qdrant_client.http import exceptions as qdrant_exceptions

# Fix for sentence-transformers compatibility
import warnings
warnings.filterwarnings("ignore", category=FutureWarning)

try:
    from sentence_transformers import SentenceTransformer
except ImportError:
    # Monkey patch for compatibility
    import huggingface_hub
    if not hasattr(huggingface_hub, 'cached_download'):
        huggingface_hub.cached_download = lambda *args, **kwargs: None
    from sentence_transformers import SentenceTransformer

# Utilities
from tqdm import tqdm
import numpy as np

# =====================
# Configuration
# =====================
BASE_DIR = Path(__file__).parent.parent
DATA_DIR = BASE_DIR / "data" / "test"
SQLITE_PATH = BASE_DIR / "data" / "index_metadata.db"
LOGS_DIR = BASE_DIR / "logs"
TEMP_DIR = BASE_DIR / "temp"

QDRANT_URL = "http://localhost:6333"
COLLECTION_NAME = "optimAIze-index"

# Embedding configuration
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
EMBEDDING_DIM = 384

# Chunking configuration
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50
MIN_CHUNK_LENGTH = 100

# Create directories
for dir_path in [DATA_DIR, LOGS_DIR, TEMP_DIR]:
    dir_path.mkdir(parents=True, exist_ok=True)

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.FileHandler(LOGS_DIR / "ingest.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# =====================
# Document Processing
# =====================
class DocumentProcessor:
    """Process various document formats including CSV and XLSX"""
    
    def __init__(self):
        self.temp_dir = TEMP_DIR
        self.temp_dir.mkdir(exist_ok=True)
        
    def extract_text(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text and metadata from document"""
        file_ext = file_path.suffix.lower()
        
        try:
            if file_ext == '.pdf':
                return self._process_pdf(file_path)
            elif file_ext == '.docx':
                return self._process_docx(file_path)
            elif file_ext == '.pptx':
                return self._process_pptx(file_path)
            elif file_ext == '.txt':
                return self._process_txt(file_path)
            elif file_ext == '.csv':
                return self._process_csv(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self._process_excel(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_ext}")
                return "", {}
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            return "", {}
    
    def _process_pdf(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Process PDF with OCR fallback"""
        try:
            doc = fitz.open(str(file_path))
            full_text = []
            metadata = {
                'page_count': len(doc),
                'has_ocr': False
            }
            
            # Extract metadata
            if doc.metadata:
                metadata.update({
                    'title': doc.metadata.get('title', ''),
                    'author': doc.metadata.get('author', ''),
                    'subject': doc.metadata.get('subject', ''),
                    'keywords': doc.metadata.get('keywords', ''),
                })
            
            for page_num, page in enumerate(doc):
                # Try embedded text first
                text = page.get_text().strip()
                
                if text and len(text) > 50:
                    full_text.append(text)
                else:
                    # Use OCR
                    logger.info(f"Using OCR for page {page_num + 1} of {file_path.name}")
                    metadata['has_ocr'] = True
                    
                    try:
                        # Render page as image
                        pix = page.get_pixmap(dpi=300)
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        
                        # Perform OCR
                        ocr_text = pytesseract.image_to_string(img, lang='eng')
                        if ocr_text.strip():
                            full_text.append(ocr_text)
                    except Exception as e:
                        logger.error(f"OCR failed for page {page_num + 1}: {e}")
            
            doc.close()
            return "\n\n".join(full_text), metadata
            
        except Exception as e:
            logger.error(f"Failed to process PDF {file_path}: {e}")
            return "", {}
    
    def _process_docx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Process DOCX file"""
        try:
            doc = DocxDocument(str(file_path))
            paragraphs = []
            
            for p in doc.paragraphs:
                if p.text and p.text.strip():
                    paragraphs.append(p.text.strip())
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        paragraphs.append(" | ".join(row_text))
            
            metadata = {
                'paragraph_count': len(paragraphs),
                'table_count': len(doc.tables)
            }
            
            # Try to get metadata
            try:
                if hasattr(doc, 'core_properties'):
                    metadata.update({
                        'author': getattr(doc.core_properties, 'author', '') or '',
                        'title': getattr(doc.core_properties, 'title', '') or '',
                    })
            except:
                pass
            
            return "\n\n".join(paragraphs), metadata
            
        except Exception as e:
            logger.error(f"Failed to process DOCX {file_path}: {e}")
            return "", {}
    
    def _process_pptx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Process PPTX file"""
        try:
            prs = Presentation(str(file_path))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if slide_text:
                    text_parts.append(f"Slide {slide_num + 1}:\n" + "\n".join(slide_text))
            
            metadata = {
                'slide_count': len(prs.slides),
            }
            
            # Try to get metadata
            try:
                if hasattr(prs, 'core_properties'):
                    metadata.update({
                        'title': getattr(prs.core_properties, 'title', '') or '',
                        'author': getattr(prs.core_properties, 'author', '') or ''
                    })
            except:
                pass
            
            return "\n\n".join(text_parts), metadata
            
        except Exception as e:
            logger.error(f"Failed to process PPTX {file_path}: {e}")
            return "", {}
    
    def _process_txt(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Process text file"""
        try:
            text = file_path.read_text(encoding='utf-8', errors='ignore')
            metadata = {'encoding': 'utf-8'}
            return text, metadata
        except Exception as e:
            logger.error(f"Failed to process TXT {file_path}: {e}")
            return "", {}
    
    def _process_csv(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Process CSV file"""
        try:
            # Read CSV with pandas
            df = pd.read_csv(file_path, encoding='utf-8', on_bad_lines='skip')
            
            # Convert to text format
            text_parts = []
            
            # Add column headers
            headers = df.columns.tolist()
            text_parts.append("CSV File with columns: " + ", ".join(str(h) for h in headers))
            
            # Add data rows (limit to prevent huge text)
            max_rows = 1000  # Adjust as needed
            for idx, row in df.head(max_rows).iterrows():
                row_text = []
                for col, val in row.items():
                    if pd.notna(val):
                        row_text.append(f"{col}: {val}")
                if row_text:
                    text_parts.append("; ".join(row_text))
            
            if len(df) > max_rows:
                text_parts.append(f"... and {len(df) - max_rows} more rows")
            
            metadata = {
                'row_count': len(df),
                'column_count': len(df.columns),
                'columns': headers
            }
            
            return "\n".join(text_parts), metadata
            
        except Exception as e:
            logger.error(f"Failed to process CSV {file_path}: {e}")
            return "", {}
    
    def _process_excel(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Process Excel file"""
        try:
            # Read Excel file
            excel_file = pd.ExcelFile(file_path)
            text_parts = []
            
            metadata = {
                'sheet_count': len(excel_file.sheet_names),
                'sheets': excel_file.sheet_names
            }
            
            # Process each sheet
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                text_parts.append(f"\n=== Sheet: {sheet_name} ===")
                
                # Add column headers
                headers = df.columns.tolist()
                text_parts.append("Columns: " + ", ".join(str(h) for h in headers))
                
                # Add data rows (limit to prevent huge text)
                max_rows = 500  # Adjust as needed
                for idx, row in df.head(max_rows).iterrows():
                    row_text = []
                    for col, val in row.items():
                        if pd.notna(val):
                            row_text.append(f"{col}: {val}")
                    if row_text:
                        text_parts.append("; ".join(row_text))
                
                if len(df) > max_rows:
                    text_parts.append(f"... and {len(df) - max_rows} more rows in sheet {sheet_name}")
            
            return "\n".join(text_parts), metadata
            
        except Exception as e:
            logger.error(f"Failed to process Excel {file_path}: {e}")
            return "", {}

# =====================
# Text Chunking
# =====================
class TextChunker:
    """Chunk text into smaller segments"""
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, 
                   overlap: int = CHUNK_OVERLAP) -> List[str]:
        """Split text into overlapping chunks"""
        if not text or len(text.strip()) < MIN_CHUNK_LENGTH:
            return []
        
        # Clean text
        text = " ".join(text.split())
        words = text.split()
        
        if len(words) <= chunk_size:
            return [text]
        
        chunks = []
        for i in range(0, len(words), chunk_size - overlap):
            chunk_words = words[i:i + chunk_size]
            chunk = " ".join(chunk_words)
            if len(chunk.strip()) >= MIN_CHUNK_LENGTH:
                chunks.append(chunk)
        
        return chunks

# =====================
# SQLite Storage
# =====================
class SQLiteStorage:
    """Handle SQLite database operations with migration support"""
    
    def __init__(self, db_path: Path):
        self.db_path = db_path
        self.init_database()
    
    def init_database(self):
        """Initialize SQLite with FTS5 and handle schema migrations"""
        conn = sqlite3.connect(self.db_path)
        c = conn.cursor()
        
        # Check if documents table exists
        c.execute("""
            SELECT name FROM sqlite_master 
            WHERE type='table' AND name='documents'
        """)
        table_exists = c.fetchone() is not None
        
        if table_exists:
            # Check if we need to migrate the schema
            c.execute("PRAGMA table_info(documents)")
            columns = [row[1] for row in c.fetchall()]
            
            if 'checksum' not in columns:
                logger.info("Migrating database schema...")
                
                # Backup existing data
                c.execute("""
                    CREATE TABLE documents_backup AS 
                    SELECT * FROM documents
                """)
                
                c.execute("""
                    CREATE TABLE chunks_backup AS 
                    SELECT * FROM chunks
                """)
                
                # Drop old tables
                c.execute("DROP TABLE IF EXISTS chunks_fts")
                c.execute("DROP TRIGGER IF EXISTS chunks_fts_insert")
                c.execute("DROP TABLE IF EXISTS chunks")
                c.execute("DROP TABLE IF EXISTS documents")
                
                # Create new schema
                self._create_schema(c)
                
                # Migrate data if possible
                try:
                    # Migrate documents
                    old_docs = c.execute("SELECT * FROM documents_backup").fetchall()
                    for doc in old_docs:
                        # Generate checksum from file path
                        file_path = doc[1] if len(doc) > 1 else ""
                        checksum = hashlib.md5(file_path.encode()).hexdigest()
                        
                        c.execute("""
                            INSERT INTO documents 
                            (document_id, file_path, file_name, file_type, file_size, checksum, ingestion_date, metadata)
                            VALUES (?, ?, ?, '.unknown', 0, ?, ?, '{}')
                        """, (doc[0], file_path, Path(file_path).name, checksum, datetime.now().isoformat()))
                    
                    # Migrate chunks
                    old_chunks = c.execute("SELECT * FROM chunks_backup").fetchall()
                    for chunk in old_chunks:
                        c.execute("""
                            INSERT INTO chunks (chunk_id, document_id, chunk_text, chunk_index)
                            VALUES (?, ?, ?, ?)
                        """, chunk[:4])
                    
                    logger.info("Migration completed successfully")
                except Exception as e:
                    logger.error(f"Migration failed: {e}")
                
                # Clean up backup tables
                c.execute("DROP TABLE documents_backup")
                c.execute("DROP TABLE chunks_backup")
        else:
            # Create fresh schema
            self._create_schema(c)
        
        conn.commit()
        conn.close()
        logger.info("SQLite database ready")
    
    def _create_schema(self, cursor):
        """Create database schema"""
        # Documents table
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS documents (
                document_id TEXT PRIMARY KEY,
                file_path TEXT UNIQUE NOT NULL,
                file_name TEXT NOT NULL,
                file_type TEXT,
                file_size INTEGER,
                checksum TEXT,
                ingestion_date TEXT,
                metadata TEXT
            )
        """)
        
        # Chunks table
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS chunks (
                chunk_id TEXT PRIMARY KEY,
                document_id TEXT NOT NULL,
                chunk_text TEXT NOT NULL,
                chunk_index INTEGER NOT NULL,
                FOREIGN KEY(document_id) REFERENCES documents(document_id)
            )
        """)
        
        # FTS5 for full-text search
        try:
            cursor.execute("""
                CREATE VIRTUAL TABLE IF NOT EXISTS chunks_fts 
                USING fts5(chunk_text, content='chunks', content_rowid='rowid')
            """)
            
            # Trigger to sync FTS
            cursor.execute("""
                CREATE TRIGGER IF NOT EXISTS chunks_fts_insert 
                AFTER INSERT ON chunks BEGIN
                    INSERT INTO chunks_fts(rowid, chunk_text) 
                    VALUES (new.rowid, new.chunk_text);
                END
            """)
        except sqlite3.OperationalError as e:
            logger.warning(f"FTS5 setup issue: {e}")
    
    def document_exists(self, checksum: str) -> bool:
        """Check if document already exists"""
        conn = sqlite3.connect(self.db_path)
        c = conn.cursor()
        result = c.execute(
            "SELECT 1 FROM documents WHERE checksum = ?", (checksum,)
        ).fetchone()
        conn.close()
        return result is not None
    
    def add_document(self, doc_info: Dict) -> str:
        """Add document to database"""
        conn = sqlite3.connect(self.db_path)
        c = conn.cursor()
        
        doc_id = str(uuid.uuid4())
        c.execute("""
            INSERT INTO documents 
            (document_id, file_path, file_name, file_type, file_size, checksum, ingestion_date, metadata)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """, (
            doc_id,
            doc_info['file_path'],
            doc_info['file_name'],
            doc_info['file_type'],
            doc_info['file_size'],
            doc_info['checksum'],
            datetime.now().isoformat(),
            json.dumps(doc_info.get('metadata', {}))
        ))
        
        conn.commit()
        conn.close()
        return doc_id
    
    def add_chunks(self, doc_id: str, chunks: List[str]) -> List[str]:
        """Add chunks to database"""
        conn = sqlite3.connect(self.db_path)
        c = conn.cursor()
        
        chunk_ids = []
        for i, chunk in enumerate(chunks):
            chunk_id = f"{doc_id}_{i}"
            c.execute("""
                INSERT INTO chunks (chunk_id, document_id, chunk_text, chunk_index)
                VALUES (?, ?, ?, ?)
            """, (chunk_id, doc_id, chunk, i))
            chunk_ids.append(chunk_id)
        
        conn.commit()
        conn.close()
        return chunk_ids

# =====================
# Main Pipeline
# =====================
class IngestionPipeline:
    """Main ingestion pipeline"""
    
    def __init__(self):
        logger.info("Initializing ingestion pipeline...")
        
        # Initialize components
        self.doc_processor = DocumentProcessor()
        self.chunker = TextChunker()
        self.storage = SQLiteStorage(SQLITE_PATH)
        
        # Initialize embedding model
        logger.info(f"Loading embedding model: {EMBEDDING_MODEL}")
        try:
            self.embed_model = SentenceTransformer(EMBEDDING_MODEL)
            logger.info("Embedding model loaded successfully")
        except Exception as e:
            logger.error(f"Failed to load embedding model: {e}")
            raise
        
        # Initialize Qdrant
        logger.info("Connecting to Qdrant...")
        try:
            self.qdrant = QdrantClient(url=QDRANT_URL)
            
            # Check if Qdrant is accessible
            collections = self.qdrant.get_collections()
            logger.info(f"Connected to Qdrant. Existing collections: {len(collections.collections)}")
        except Exception as e:
            logger.error(f"Failed to connect to Qdrant: {e}")
            logger.error("Make sure Qdrant is running: docker run -p 6333:6333 qdrant/qdrant")
            raise
        
        # Handle collection setup
        self._setup_collection()
    
    def _setup_collection(self):
        """Setup Qdrant collection with proper error handling"""
        try:
            # Try to get collection info
            collection_info = self.qdrant.get_collection(COLLECTION_NAME)
            vectors_count = collection_info.vectors_count if hasattr(collection_info, 'vectors_count') else 0
            logger.info(f"Collection '{COLLECTION_NAME}' exists with {vectors_count} vectors")
            
            # Check if we should recreate
            if vectors_count > 0:
                response = input(f"\nCollection '{COLLECTION_NAME}' already has {vectors_count} vectors. Delete and recreate? (y/n): ")
                if response.lower() == 'y':
                    logger.info("Deleting existing collection...")
                    self.qdrant.delete_collection(COLLECTION_NAME)
                    time.sleep(1)
                    self._create_collection()
                else:
                    logger.info("Keeping existing collection")
            else:
                # Empty collection, can use as is
                logger.info("Using existing empty collection")
                
        except Exception as e:
            # Collection doesn't exist, create it
            logger.info(f"Collection '{COLLECTION_NAME}' does not exist, creating it...")
            self._create_collection()
    
    def _create_collection(self):
        """Create new collection"""
        max_retries = 3
        for attempt in range(max_retries):
            try:
                self.qdrant.create_collection(
                    collection_name=COLLECTION_NAME,
                    vectors_config=VectorParams(size=EMBEDDING_DIM, distance=Distance.COSINE)
                )
                logger.info(f"Created collection '{COLLECTION_NAME}'")
                break
            except Exception as e:
                if attempt < max_retries - 1:
                    logger.warning(f"Failed to create collection (attempt {attempt + 1}): {e}")
                    time.sleep(2)
                else:
                    logger.error(f"Failed to create collection after {max_retries} attempts: {e}")
                    raise
    
    def calculate_checksum(self, file_path: Path) -> str:
        """Calculate file checksum"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    def process_file(self, file_path: Path) -> bool:
        """Process a single file"""
        try:
            logger.info(f"Processing: {file_path}")
            
            # Calculate checksum
            checksum = self.calculate_checksum(file_path)
            
            # Check if already processed
            if self.storage.document_exists(checksum):
                logger.info(f"Skipping (already processed): {file_path}")
                return True
            
            # Extract text and metadata
            text, metadata = self.doc_processor.extract_text(file_path)
            
            if not text or len(text.strip()) < MIN_CHUNK_LENGTH:
                logger.warning(f"No meaningful content extracted from {file_path}")
                return False
            
            logger.info(f"Extracted {len(text)} characters from {file_path}")
            
            # Chunk text
            chunks = self.chunker.chunk_text(text)
            logger.info(f"Created {len(chunks)} chunks")
            
            if not chunks:
                logger.warning(f"No chunks created for {file_path}")
                return False
            
            # Store in SQLite
            doc_info = {
                'file_path': str(file_path),
                'file_name': file_path.name,
                'file_type': file_path.suffix.lower(),
                'file_size': file_path.stat().st_size,
                'checksum': checksum,
                'metadata': metadata
            }
            
            doc_id = self.storage.add_document(doc_info)
            chunk_ids = self.storage.add_chunks(doc_id, chunks)
            
            # Generate embeddings in batches
            logger.info("Generating embeddings...")
            batch_size = 32
            all_embeddings = []
            
            for i in range(0, len(chunks), batch_size):
                batch = chunks[i:i + batch_size]
                try:
                    batch_embeddings = self.embed_model.encode(batch, show_progress_bar=False)
                    all_embeddings.extend(batch_embeddings)
                except Exception as e:
                    logger.error(f"Failed to generate embeddings for batch: {e}")
                    return False
            
            logger.info(f"Generated {len(all_embeddings)} embeddings")
            
            # Create Qdrant points in batches
            points = []
            for chunk_id, chunk, embedding in zip(chunk_ids, chunks, all_embeddings):
                point = PointStruct(
                    id=chunk_id,
                    vector=embedding.tolist(),
                    payload={
                        'document_id': doc_id,
                        'file_name': file_path.name,
                        'file_path': str(file_path),
                        'chunk_text': chunk[:500],  # Truncate for payload
                        'chunk_index': chunk_ids.index(chunk_id),
                        'file_type': file_path.suffix.lower(),
                        **metadata
                    }
                )
                points.append(point)
            
            # Upsert to Qdrant in batches
            batch_size = 100
            for i in range(0, len(points), batch_size):
                batch = points[i:i + batch_size]
                try:
                    self.qdrant.upsert(
                        collection_name=COLLECTION_NAME,
                        points=batch
                    )
                except Exception as e:
                    logger.error(f"Failed to upsert batch to Qdrant: {e}")
                    return False
            
            logger.info(f"Successfully indexed {file_path} ({len(points)} chunks)")
            return True
            
        except Exception as e:
            logger.error(f"Failed to process {file_path}: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def ingest_directory(self, directory: Path):
        """Process all files in directory"""
        # Find all supported files
        patterns = ['*.pdf', '*.docx', '*.pptx', '*.txt', '*.csv', '*.xlsx', '*.xls']
        all_files = []
        
        for pattern in patterns:
            all_files.extend(directory.glob(pattern))
        
        # Skip database and temporary files
        all_files = [f for f in all_files if not f.name.startswith('.') and f.name != "index_metadata.db"]
        
        if not all_files:
            logger.warning(f"No supported files found in {directory}")
            logger.info("Supported formats: .pdf, .docx, .pptx, .txt, .csv, .xlsx, .xls")
            return
        
        logger.info(f"Found {len(all_files)} files to process:")
        for f in all_files:
            logger.info(f"  - {f.name} ({f.suffix})")
        
        success_count = 0
        failed_files = []
        
        for file_path in tqdm(all_files, desc="Processing files"):
            if self.process_file(file_path):
                success_count += 1
            else:
                failed_files.append(file_path)
        
        logger.info(f"\nSuccessfully processed {success_count}/{len(all_files)} files")
        
        if failed_files:
            logger.warning("Failed files:")
            for f in failed_files:
                logger.warning(f"  - {f}")
        
        # Verify results
        self.verify_ingestion()
    
    def verify_ingestion(self):
        """Verify ingestion results"""
        logger.info("\n" + "=" * 50)
        logger.info("VERIFICATION RESULTS")
        logger.info("=" * 50)
        
        # Check Qdrant
        try:
            collection_info = self.qdrant.get_collection(COLLECTION_NAME)
            # Handle different response formats
            if hasattr(collection_info, 'vectors_count'):
                vector_count = collection_info.vectors_count
            elif hasattr(collection_info, 'points_count'):
                vector_count = collection_info.points_count
            else:
                # Try to count vectors manually
                vector_count = len(self.qdrant.scroll(
                    collection_name=COLLECTION_NAME,
                    limit=1,
                    with_payload=False,
                    with_vectors=False
                )[0])
            
            logger.info(f"✅ Qdrant vectors: {vector_count}")
        except Exception as e:
            logger.error(f"❌ Failed to get Qdrant info: {e}")
        
        # Check SQLite
        try:
            conn = sqlite3.connect(SQLITE_PATH)
            c = conn.cursor()
            
            doc_count = c.execute("SELECT COUNT(*) FROM documents").fetchone()[0]
            chunk_count = c.execute("SELECT COUNT(*) FROM chunks").fetchone()[0]
            
            logger.info(f"✅ SQLite documents: {doc_count}")
            logger.info(f"✅ SQLite chunks: {chunk_count}")
            
            # Show document types
            try:
                doc_types = c.execute("""
                    SELECT file_type, COUNT(*) as count 
                    FROM documents 
                    GROUP BY file_type
                    ORDER BY count DESC
                """).fetchall()
                
                if doc_types:
                    logger.info("\nDocument types:")
                    for file_type, count in doc_types:
                        logger.info(f"  - {file_type}: {count} files")
            except Exception as e:
                logger.warning(f"Could not fetch document types: {e}")
            
            # Test FTS5 search
            try:
                test_results = c.execute(
                    "SELECT COUNT(*) FROM chunks_fts WHERE chunk_text MATCH 'test'"
                ).fetchone()[0]
                logger.info(f"✅ FTS5 search test: {test_results} results for 'test'")
            except Exception as e:
                logger.warning(f"⚠️  FTS5 search not available: {e}")
            
            conn.close()
        except Exception as e:
            logger.error(f"❌ Failed to check SQLite: {e}")
        
        # Test vector search
        try:
            test_query = "sample test query"
            test_embedding = self.embed_model.encode([test_query])[0]
            
            test_results = self.qdrant.search(
                collection_name=COLLECTION_NAME,
                query_vector=test_embedding.tolist(),
                limit=3
            )
            logger.info(f"\n✅ Vector search test: {len(test_results)} results")
            
            if test_results:
                logger.info("\nSample results for query 'sample test query':")
                for i, result in enumerate(test_results[:3], 1):
                    logger.info(f"\n{i}. File: {result.payload.get('file_name', 'Unknown')}")
                    logger.info(f"   Type: {result.payload.get('file_type', 'Unknown')}")
                    logger.info(f"   Score: {result.score:.4f}")
                    preview = result.payload.get('chunk_text', '')[:100] + "..."
                    logger.info(f"   Preview: {preview}")
        except Exception as e:
            logger.error(f"❌ Vector search test failed: {e}")

# =====================
# Main execution
# =====================
def main():
    logger.info("=" * 50)
    logger.info("Document Ingestion Pipeline")
    logger.info("=" * 50)
    
    # Check if data directory exists
    if not DATA_DIR.exists():
        logger.error(f"Data directory does not exist: {DATA_DIR}")
        DATA_DIR.mkdir(parents=True, exist_ok=True)
        logger.info(f"Created data directory: {DATA_DIR}")
        
        # Create a sample file
        sample_file = DATA_DIR / "sample.txt"
        sample_file.write_text("""
Sample Document for Testing

This is a sample document to test the ingestion pipeline.
It contains multiple paragraphs to test chunking.

The system should process this document and create embeddings.
Each chunk will be stored in both Qdrant and SQLite.

Features being tested:
- Text extraction from various formats (PDF, DOCX, PPTX, TXT, CSV, XLSX)
- Chunking with overlap
- Embedding generation using sentence-transformers
- Vector storage in Qdrant
- Metadata storage in SQLite with FTS5
- Full-text search capabilities
- Deduplication using checksums

This pipeline supports:
1. PDF files with OCR for scanned pages
2. Word documents with table extraction
3. PowerPoint presentations
4. Plain text files
5. CSV files with structured data
6. Excel files with multiple sheets

The system is designed to be robust and handle edge cases gracefully.
        """)
        logger.info(f"Created sample file: {sample_file}")
    
    # Check for Tesseract installation
    try:
        import subprocess
        result = subprocess.run(['tesseract', '--version'], capture_output=True, text=True)
        if result.returncode == 0:
            logger.info(f"Tesseract is installed: {result.stdout.split()[1]}")
        else:
            logger.warning("Tesseract not found - OCR will not work for scanned PDFs")
            logger.info("Install with: brew install tesseract (macOS) or apt-get install tesseract-ocr (Linux)")
    except FileNotFoundError:
        logger.warning("Tesseract not found - OCR will not work for scanned PDFs")
    
    try:
        pipeline = IngestionPipeline()
        pipeline.ingest_directory(DATA_DIR)
        logger.info("\n✅ Ingestion complete!")
        
        # Print summary
        logger.info("\n" + "=" * 50)
        logger.info("SUMMARY")
        logger.info("=" * 50)
        logger.info(f"Data directory: {DATA_DIR}")
        logger.info(f"SQLite database: {SQLITE_PATH}")
        logger.info(f"Qdrant URL: {QDRANT_URL}")
        logger.info(f"Collection: {COLLECTION_NAME}")
        logger.info(f"Embedding model: {EMBEDDING_MODEL}")
        logger.info(f"Chunk size: {CHUNK_SIZE} words")
        logger.info(f"Chunk overlap: {CHUNK_OVERLAP} words")
        
    except Exception as e:
        logger.error(f"\n❌ Pipeline failed: {e}")
        import traceback
        traceback.print_exc()
        
        # Provide troubleshooting tips
        logger.info("\n" + "=" * 50)
        logger.info("TROUBLESHOOTING")
        logger.info("=" * 50)
        
        if "Qdrant" in str(e):
            logger.info("1. Make sure Qdrant is running:")
            logger.info("   docker run -p 6333:6333 qdrant/qdrant")
            logger.info("\n2. Check if port 6333 is already in use:")
            logger.info("   lsof -i :6333")
            logger.info("\n3. Try restarting the Qdrant container:")
            logger.info("   docker restart <container_id>")
        
        if "sentence_transformers" in str(e) or "huggingface" in str(e):
            logger.info("1. Try reinstalling sentence-transformers:")
            logger.info("   pip uninstall sentence-transformers huggingface-hub -y")
            logger.info("   pip install sentence-transformers==2.5.1 huggingface-hub==0.20.3")
        
        if "SQLite" in str(e):
            logger.info("1. Try deleting the existing database:")
            logger.info(f"   rm {SQLITE_PATH}")
            logger.info("\n2. The pipeline will create a fresh database on next run")

if __name__ == "__main__":
    main()
    